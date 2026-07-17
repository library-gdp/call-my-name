#!/usr/bin/env python3
"""Extract user-authored text from supported document formats."""

from __future__ import annotations

import argparse
import os
import sys
import tempfile
from pathlib import Path
from typing import Iterable


TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".xml",
    ".html",
    ".htm",
}


class UnsupportedFormatError(Exception):
    """Raised when no bundled parser supports the input suffix."""


class DocumentParseError(Exception):
    """Raised when a supported document cannot yield usable text."""


def _clean_lines(lines: Iterable[str]) -> list[str]:
    cleaned: list[str] = []
    previous_blank = False
    for raw_line in lines:
        line = " ".join(raw_line.replace("\x00", "").split())
        if not line:
            if cleaned and not previous_blank:
                cleaned.append("")
            previous_blank = True
            continue
        cleaned.append(line)
        previous_blank = False
    while cleaned and not cleaned[-1]:
        cleaned.pop()
    return cleaned


def _read_text(path: Path) -> str:
    data = path.read_bytes()
    if b"\x00" in data:
        raise DocumentParseError("text file contains NUL bytes")
    for encoding in ("utf-8-sig", "utf-8", "cp949"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentParseError("text encoding is not UTF-8 or CP949")


def _parse_pdf(path: Path) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted and reader.decrypt("") == 0:
        raise DocumentParseError("PDF is encrypted")

    lines: list[str] = []
    extracted_pages = 0
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            extracted_pages += 1
            lines.extend(text.splitlines())
            lines.append("")
    if extracted_pages == 0:
        raise DocumentParseError("PDF has no extractable text; OCR is not bundled")
    return lines


def _table_lines(table: object) -> list[str]:
    lines: list[str] = []
    for row in table.rows:  # type: ignore[attr-defined]
        values = [cell.text.strip() for cell in row.cells]
        if any(values):
            lines.append("\t".join(values))
    return lines


def _parse_docx(path: Path) -> list[str]:
    from docx import Document
    from docx.table import Table

    document = Document(str(path))
    lines: list[str] = []

    boundary_text: list[str] = []
    for section in document.sections:
        for container in (section.header, section.footer):
            boundary_text.extend(paragraph.text for paragraph in container.paragraphs)
            for table in container.tables:
                boundary_text.extend(_table_lines(table))
    lines.extend(dict.fromkeys(text for text in boundary_text if text.strip()))
    if lines:
        lines.append("")

    if hasattr(document, "iter_inner_content"):
        for item in document.iter_inner_content():
            if isinstance(item, Table):
                lines.extend(_table_lines(item))
            else:
                lines.append(item.text)
    else:
        lines.extend(paragraph.text for paragraph in document.paragraphs)
        for table in document.tables:
            lines.extend(_table_lines(table))
    return lines


def _format_cell(value: object) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        try:
            return value.isoformat()  # type: ignore[no-any-return]
        except TypeError:
            pass
    return str(value).strip()


def _parse_xlsx(path: Path) -> list[str]:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=path, read_only=True, data_only=False)
    lines: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            lines.append(worksheet.title)
            for row in worksheet.iter_rows(values_only=True):
                values = [_format_cell(value) for value in row]
                if any(values):
                    lines.append("\t".join(values))
            lines.append("")
    finally:
        workbook.close()
    return lines


def _shape_lines(shape: object) -> list[str]:
    lines: list[str] = []
    if getattr(shape, "has_text_frame", False):
        lines.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
    if getattr(shape, "has_table", False):
        lines.extend(_table_lines(shape.table))
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            lines.extend(_shape_lines(child))
    return lines


def _parse_pptx(path: Path) -> list[str]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            lines.extend(_shape_lines(shape))
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                lines.extend(paragraph.text for paragraph in notes_frame.paragraphs)
        lines.append("")
    return lines


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        lines = _read_text(path).splitlines()
    elif suffix == ".pdf":
        lines = _parse_pdf(path)
    elif suffix == ".docx":
        lines = _parse_docx(path)
    elif suffix == ".xlsx":
        lines = _parse_xlsx(path)
    elif suffix == ".pptx":
        lines = _parse_pptx(path)
    else:
        raise UnsupportedFormatError(f"unsupported file extension: {suffix or '<none>'}")

    cleaned = _clean_lines(lines)
    if not any(cleaned):
        raise DocumentParseError("document contains no usable text")
    return "\n".join(cleaned) + "\n"


def _write_atomic(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    handle = tempfile.NamedTemporaryFile(
        mode="w",
        encoding="utf-8",
        newline="\n",
        dir=path.parent,
        prefix=f".{path.name}.",
        suffix=".tmp",
        delete=False,
    )
    temporary_path = Path(handle.name)
    try:
        with handle:
            handle.write(content)
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temporary_path, path)
    except BaseException:
        temporary_path.unlink(missing_ok=True)
        raise


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, type=Path, help="input document")
    parser.add_argument("--output", required=True, type=Path, help="UTF-8 text output")
    return parser.parse_args()


def main() -> int:
    args = _parse_args()
    try:
        if args.input.is_symlink():
            raise DocumentParseError("input must not be a symlink")
        input_path = args.input.resolve(strict=True)
        if not input_path.is_file():
            raise DocumentParseError("input must be a regular file")
        _write_atomic(args.output, extract(input_path))
    except UnsupportedFormatError as error:
        print(str(error), file=sys.stderr)
        return 2
    except (DocumentParseError, OSError, ValueError) as error:
        print(f"failed to parse {args.input}: {error}", file=sys.stderr)
        return 3
    except Exception as error:
        print(f"failed to parse {args.input}: {error}", file=sys.stderr)
        return 3
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
